import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 슬라이드 크기를 16:9 와이드스크린으로 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 공통 컬러 정의
    COLOR_PRIMARY_DARK = RGBColor(34, 34, 59)     # 진한 남보라
    COLOR_SECONDARY_DARK = RGBColor(74, 78, 105)   # 회보라
    COLOR_PRIMARY_LIGHT = RGBColor(245, 246, 250)  # 아주 밝은 회색 (배경용)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(40, 40, 40)
    COLOR_TEXT_LIGHT = RGBColor(220, 220, 220)
    COLOR_ACCENT = RGBColor(180, 50, 120)          # 핫핑크 Accent
    COLOR_PURPLE = RGBColor(108, 92, 231)          # 보라색 테마 Accent
    COLOR_CYAN = RGBColor(0, 184, 148)             # 청록색 Accent (새 슬라이드용)

    # ==========================================
    # 슬라이드 1: 표지 (Dark 테마)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # 빈 슬라이드
    slide1 = prs.slides.add_slide(slide_layout)
    
    # 어두운 배경 추가
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    bg.line.fill.background()

    # 상단 장식 라인
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # 메인 타이틀 박스
    title_box = slide1.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "AI 기반 H-Beam 도면 적산 검증 시스템"
    p1.font.name = "Malgun Gothic"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "CAD 도면 파싱과 Gemini AI의 시너지를 통한 철골 적산 혁신"
    p2.font.name = "Malgun Gothic"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_TEXT_LIGHT

    # 하단 설명 정보
    info_box = slide1.shapes.add_textbox(Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.5))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "주요 기술: ezdxf 파싱 엔진  |  규칙/AI 하이브리드 탐지  |  예외 복구용 수동 시트 분할"
    p_info.font.name = "Malgun Gothic"
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = COLOR_SECONDARY_DARK


    # ==========================================
    # 슬라이드 2: 도면 분석의 특징 및 한계점 (Light 테마)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
    bg2.line.fill.background()

    # 슬라이드 제목
    title_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf2 = title_box2.text_frame
    p_t2 = tf2.paragraphs[0]
    p_t2.text = "CAD 도면 분석의 특징과 자동화의 어려움"
    p_t2.font.name = "Malgun Gothic"
    p_t2.font.size = Pt(28)
    p_t2.font.bold = True
    p_t2.font.color.rgb = COLOR_PRIMARY_DARK

    # 왼쪽 카드: 도면 데이터의 특징
    card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_WHITE
    card_left.line.color.rgb = RGBColor(220, 220, 220)
    card_left.line.width = Pt(1)

    tf_left = card_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0.3)
    tf_left.margin_right = Inches(0.3)
    tf_left.margin_top = Inches(0.3)

    p_l1 = tf_left.paragraphs[0]
    p_l1.text = "도면 데이터의 고유 특징"
    p_l1.font.name = "Malgun Gothic"
    p_l1.font.size = Pt(20)
    p_l1.font.bold = True
    p_l1.font.color.rgb = COLOR_PURPLE
    p_l1.space_after = Pt(18)

    bullets_left = [
        "선(Line), 원(Circle), 텍스트(Text) 등의 기하학적 요소가 복잡하게 얽혀 있는 고유의 비정형 데이터 구조",
        "동일 부재라 할지라도 도면 작성자마다 표현 방식, 폰트 및 축척이 다르며 정형화된 표준 가이드 부재",
        "부재 목록(참조 시트) 정보와 평면도 내 도면 요소 간의 관계 유기적 파악 필요"
    ]
    for b in bullets_left:
        p_b = tf_left.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_after = Pt(12)

    # 오른쪽 카드: 자동화의 어려움 (한계)
    card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_WHITE
    card_right.line.color.rgb = RGBColor(220, 220, 220)
    card_right.line.width = Pt(1)

    tf_right = card_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.3)
    tf_right.margin_right = Inches(0.3)
    tf_right.margin_top = Inches(0.3)

    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "자동 분석의 어려움 및 한계점"
    p_r1.font.name = "Malgun Gothic"
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = COLOR_ACCENT
    p_r1.space_after = Pt(18)

    bullets_right = [
        "규칙 기반의 한계: 불규칙한 텍스트 폰트 깨짐이나 예외적인 간섭 발생 시 단순 룰 알고리즘은 인식 누락 발생",
        "배경 노이즈 간섭: 치수선, 그리드 해치 패턴 등 다양한 배경 도면 정보가 뒤섞여 순수 부재 정보 필터링 난해",
        "다중 시트 분할 및 좌표 오차: 하나의 큰 파일 내 여러 상세 도면들이 나열되어 축척 및 좌표계 매핑 오차 통제 필요"
    ]
    for b in bullets_right:
        p_b = tf_right.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_after = Pt(12)


    # ==========================================
    # [NEW] 슬라이드 3: 예외 도면 대응을 위한 수동 시트 분할 (Light 테마)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
    bg3.line.fill.background()

    # 슬라이드 제목
    title_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf3 = title_box3.text_frame
    p_t3 = tf3.paragraphs[0]
    p_t3.text = "자동 검증 한계 극복을 위한 수동 시트 분할"
    p_t3.font.name = "Malgun Gothic"
    p_t3.font.size = Pt(28)
    p_t3.font.bold = True
    p_t3.font.color.rgb = COLOR_PRIMARY_DARK

    # 왼쪽 카드: 수동 시트 분할의 도입 배경
    card3_left = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card3_left.fill.solid()
    card3_left.fill.fore_color.rgb = COLOR_WHITE
    card3_left.line.color.rgb = RGBColor(220, 220, 220)
    card3_left.line.width = Pt(1)

    tf3_left = card3_left.text_frame
    tf3_left.word_wrap = True
    tf3_left.margin_left = Inches(0.3)
    tf3_left.margin_right = Inches(0.3)
    tf3_left.margin_top = Inches(0.3)

    p3_l1 = tf3_left.paragraphs[0]
    p3_l1.text = "수동 분할 도입 배경 및 목적"
    p3_l1.font.name = "Malgun Gothic"
    p3_l1.font.size = Pt(20)
    p3_l1.font.bold = True
    p3_l1.font.color.rgb = COLOR_CYAN
    p3_l1.space_after = Pt(18)

    bullets3_left = [
        "자동 분할 실패 대응: 도각(도면 테두리)선이 손상되었거나 인쇄용 템플릿이 없어 자동 시트 분할 알고리즘이 동작하지 않는 예외 도면 대응",
        "정밀 개별 검증: 사용자가 전체 대형 도면 내에서 검토하고자 하는 특정 상세 구획만 수동으로 분할하여 분석 집중도 극대화",
        "오작동 방지 및 휴먼 에러 보완: 시스템 자동 분석이 놓치거나 잘못 판정한 레이아웃 구조를 사용자가 수동으로 직접 교정"
    ]
    for b in bullets3_left:
        p_b = tf3_left.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_after = Pt(12)

    # 오른쪽 카드: 상세 작동 방식 및 기능
    card3_right = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    card3_right.fill.solid()
    card3_right.fill.fore_color.rgb = COLOR_WHITE
    card3_right.line.color.rgb = RGBColor(220, 220, 220)
    card3_right.line.width = Pt(1)

    tf3_right = card3_right.text_frame
    tf3_right.word_wrap = True
    tf3_right.margin_left = Inches(0.3)
    tf3_right.margin_right = Inches(0.3)
    tf3_right.margin_top = Inches(0.3)

    p3_r1 = tf3_right.paragraphs[0]
    p3_r1.text = "작동 방식 및 주요 편의 기능"
    p3_r1.font.name = "Malgun Gothic"
    p3_r1.font.size = Pt(20)
    p3_r1.font.bold = True
    p3_r1.font.color.rgb = COLOR_PRIMARY_DARK
    p3_r1.space_after = Pt(18)

    bullets3_right = [
        "마우스 드래그 영역 지정: 사용자가 전체 도면 뷰어 상에서 원하는 영역을 마우스로 드래그하여 직관적으로 분할 영역 지정",
        "독립 메타데이터 연동: 분할된 각 구역에 도면 번호 및 시트 이름을 자유롭게 지정하여 개별 적산 리스트 및 마커와 독립 연동",
        "안전한 편집 이력 데이터 복원: 수동 시트 영역을 삭제하거나 수정하더라도 기존 적산 내역과 마커 정보가 다른 시트로 귀속될 수 있게 예외 방어 설계"
    ]
    for b in bullets3_right:
        p_b = tf3_right.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_after = Pt(12)


    # ==========================================
    # 슬라이드 4: 규칙과 AI의 하이브리드 검출 및 분석 (Light 테마)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
    bg4.line.fill.background()

    # 슬라이드 제목
    title_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf4 = title_box4.text_frame
    p_t4 = tf4.paragraphs[0]
    p_t4.text = "규칙과 AI의 하이브리드 검출 및 분석"
    p_t4.font.name = "Malgun Gothic"
    p_t4.font.size = Pt(28)
    p_t4.font.bold = True
    p_t4.font.color.rgb = COLOR_PRIMARY_DARK

    # 3개 열의 카드 배치
    col_width = Inches(3.64)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)
    height_pos = Inches(4.7)

    features = [
        {
            "num": "01",
            "title": "CAD 도면 파싱 & 분할",
            "desc": "ezdxf 파싱 엔진을 통해 CAD 도면 내 벡터 데이터와 텍스트를 초고속으로 추출합니다. 복잡한 다중 도면 시트를 개별 작업 영역으로 자동/수동 분할하여 독립 관리합니다."
        },
        {
            "num": "02",
            "title": "룰 기반 자동 감지",
            "desc": "사전 정의된 규칙 알고리즘을 기반으로 기둥 원형 마커 및 H-Beam 배선을 정밀 감지합니다. 부재의 고유 좌표와 실측 길이를 실시간 측정하여 목록화합니다."
        },
        {
            "num": "03",
            "title": "Gemini AI Helper 보정",
            "desc": "기존 알고리즘이 간섭과 불규칙성으로 인해 찾아내지 못한 기둥 및 보를 Gemini 2.5-flash AI가 정밀 분석 및 보정 제안하며, 신뢰도와 탐지 근거를 동시 제공합니다."
        }
    ]

    for i, feat in enumerate(features):
        left_pos = start_left + i * (col_width + col_gap)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, height_pos)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(225, 225, 230)
        card.line.width = Pt(1.5)

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_right = Inches(0.25)
        tf_c.margin_top = Inches(0.3)

        # 번호
        p_num = tf_c.paragraphs[0]
        p_num.text = feat["num"]
        p_num.font.name = "Malgun Gothic"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_PURPLE
        p_num.space_after = Pt(10)

        # 제목
        p_title = tf_c.add_paragraph()
        p_title.text = feat["title"]
        p_title.font.name = "Malgun Gothic"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_DARK
        p_title.space_after = Pt(14)

        # 상세
        p_desc = tf_c.add_paragraph()
        p_desc.text = feat["desc"]
        p_desc.font.name = "Malgun Gothic"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_SECONDARY_DARK
        p_desc.line_spacing = 1.3


    # ==========================================
    # 슬라이드 5: 사용자 편의성 및 실무 기대 효과 (Light 테마)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
    bg5.line.fill.background()

    # 슬라이드 제목
    title_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf5 = title_box5.text_frame
    p_t5 = tf5.paragraphs[0]
    p_t5.text = "직관적인 검증 환경 및 최종 리포트 자동화"
    p_t5.font.name = "Malgun Gothic"
    p_t5.font.size = Pt(28)
    p_t5.font.bold = True
    p_t5.font.color.rgb = COLOR_PRIMARY_DARK

    # 3개 가로형 박스 배치
    row_height = Inches(1.3)
    row_gap = Inches(0.25)
    row_top_start = Inches(1.7)
    row_width = Inches(11.733)

    benefits = [
        {
            "title": "동적 포커싱 & 핫핑크 하이라이트",
            "desc": "AI가 추천한 항목 카드를 클릭하는 즉시 도면 상의 해당 위치로 화면이 자동 정렬되고, 1.8배 확대(Zoom-in) 및 깜빡이는 핫핑크색 이중 포커스 링으로 탐지 대상을 직관적으로 식별할 수 있습니다."
        },
        {
            "title": "부호별 통계 요약 대시보드 및 실시간 동기화",
            "desc": "인터페이스 상단에 AI 총 제안량 및 부호별 집계 통계를 시각화하고, 사용자가 승인 또는 거절 버튼을 누르는 순간 실시간으로 대시보드 통계치가 갱신되어 최신 적산 상태를 한눈에 모니터링합니다."
        },
        {
            "title": "참조 분석기 탑재 및 Excel 내보내기",
            "desc": "높이 기준 정보를 파싱하여 연동하는 참조 분석기를 활용하고, 적산이 완료된 모든 정보는 최종적으로 정갈하고 가독성이 극대화된 표준 엑셀 리포트(.xlsx) 서식으로 자동 저장 및 다운로드합니다."
        }
    ]

    for i, ben in enumerate(benefits):
        current_top = row_top_start + i * (row_height + row_gap)
        row_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), current_top, row_width, row_height)
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = COLOR_WHITE
        row_box.line.color.rgb = RGBColor(230, 230, 235)
        row_box.line.width = Pt(1.5)

        tf_b = row_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = Inches(0.3)
        tf_b.margin_right = Inches(0.3)
        tf_b.margin_top = Inches(0.15)
        tf_b.margin_bottom = Inches(0.15)

        # 제목
        p_bt = tf_b.paragraphs[0]
        p_bt.text = ben["title"]
        p_bt.font.name = "Malgun Gothic"
        p_bt.font.size = Pt(16)
        p_bt.font.bold = True
        p_bt.font.color.rgb = COLOR_PURPLE
        p_bt.space_after = Pt(4)

        # 내용
        p_bd = tf_b.add_paragraph()
        p_bd.text = ben["desc"]
        p_bd.font.name = "Malgun Gothic"
        p_bd.font.size = Pt(12.5)
        p_bd.font.color.rgb = COLOR_TEXT_DARK
        p_bd.line_spacing = 1.25

    # 저장
    output_dir = "output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    file_path = os.path.join(output_dir, "H-Beam_Takeoff_System_Report.pptx")
    # 이미 파일이 열려 있어서 생기는 권한 에러 방지
    counter = 1
    base_name = "H-Beam_Takeoff_System_Report"
    while True:
        try:
            prs.save(file_path)
            break
        except PermissionError:
            file_path = os.path.join(output_dir, f"{base_name}_{counter}.pptx")
            counter += 1
    print(f"Presentation saved successfully at: {file_path}")

if __name__ == "__main__":
    create_deck()
